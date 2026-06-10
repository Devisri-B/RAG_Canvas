"""
Chunk metadata schema (what every chunk carries):
    week         module name, e.g. "week1"      (None in --all-files mode)
    source_file  original filename
    file_id      Canvas file id (stable handle for re-download / citation)
    loc_type     "page" (PDF) or "slide" (PPTX)
    loc          1-based page/slide number       <- powers "source reference"
    chunk_index  position of this chunk within that page/slide
    course_id    Canvas course id
    n_tokens     token count (cl100k_base)

Output: one JSON object per line in data/chunks.jsonl:
    {"id": "...", "text": "...", "metadata": {...}}

Usage:
    python src/ingest.py --course-id 2
    python src/ingest.py --course-id 2 --all-files     # ignore modules
    python src/ingest.py --course-id 2 --ocr           # also OCR image-locked content
"""

from __future__ import annotations

import argparse
import json
import re
import tempfile
import warnings
from pathlib import Path

import fitz  # PyMuPDF
import requests
import tiktoken
from canvasapi import Canvas
from dotenv import dotenv_values
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# canvasapi warns on every HTTP (non-HTTPS) request; our local Canvas is HTTP.
warnings.filterwarnings("ignore", message="Canvas may respond unexpectedly")

REPO_ROOT = Path(__file__).resolve().parent.parent
ENV_PATH = REPO_ROOT / ".env"


CHUNK_TOKENS = 900
CHUNK_OVERLAP = 120
MIN_OCR_IMG_BYTES = 15_000  # skip icons/bullets/logos when OCR-ing pictures

_enc = tiktoken.get_encoding("cl100k_base")
_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
    encoding_name="cl100k_base",
    chunk_size=CHUNK_TOKENS,
    chunk_overlap=CHUNK_OVERLAP,
)


# --------------------------------------------------------------------------- #
# Canvas connection
# --------------------------------------------------------------------------- #
def get_canvas() -> Canvas:
    cfg = dotenv_values(ENV_PATH)
    url = (cfg.get("CANVAS_BASE_URL") or "").strip()
    token = (cfg.get("CANVAS_API_TOKEN") or "").strip()
    if not url or not token:
        raise SystemExit(f"CANVAS_BASE_URL / CANVAS_API_TOKEN missing in {ENV_PATH}")
    return Canvas(url, token)


# --------------------------------------------------------------------------- #
# Enumerate the files to ingest (with their week label)
# --------------------------------------------------------------------------- #
SUPPORTED = (".pdf", ".pptx")


def iter_content_units(course, all_files: bool = False):
    """Yield (week, file_obj) for every PDF/PPTX to ingest.

    Default: walk Modules and take File items -> week = module name.
    --all-files: take every supported file on the course -> week = None.
    """
    if all_files:
        for f in course.get_files():
            if f.display_name.lower().endswith(SUPPORTED):
                yield None, f
        return

    modules = list(course.get_modules())
    if not modules:
        raise SystemExit(
            "No modules found. Organize content into Modules (e.g. 'week1'), "
            "or re-run with --all-files to ingest every file."
        )
    for m in modules:
        for item in m.get_module_items():
            if item.type != "File":
                continue
            f = course.get_file(item.content_id)
            if f.display_name.lower().endswith(SUPPORTED):
                yield m.name, f


def download(file_obj, raw_dir: Path) -> Path:
    """Download a Canvas file to raw_dir/<file_id>_<name>, cached if present."""
    raw_dir.mkdir(parents=True, exist_ok=True)
    dest = raw_dir / f"{file_obj.id}_{file_obj.display_name}"
    if not dest.exists():
        resp = requests.get(file_obj.url, timeout=60)
        resp.raise_for_status()
        dest.write_bytes(resp.content)
    return dest


# --------------------------------------------------------------------------- #
# OCR (optional, offline) — recover text locked inside images.
# --------------------------------------------------------------------------- #
_ocr_conv = None


def _ocr_image(blob: bytes, ext: str) -> str:
    """OCR one image (bytes) -> cleaned text, or '' if nothing usable."""
    global _ocr_conv
    if _ocr_conv is None:
        from docling.document_converter import DocumentConverter
        _ocr_conv = DocumentConverter()
    with tempfile.NamedTemporaryFile(suffix=f".{ext}") as tmp:
        tmp.write(blob)
        tmp.flush()
        try:
            md = _ocr_conv.convert(tmp.name).document.export_to_markdown()
        except Exception:
            return ""
    return _clean_ocr(md)


def _clean_ocr(md: str) -> str:
    """Strip docling markdown noise; drop output with too little real text.

    NOTE: this won't catch every decorative image (e.g. a stylised code-art
    photo can OCR to real-looking words). A confidence threshold via docling's
    lower-level OCR API is the robust upgrade if noise becomes a problem.
    """
    md = md.replace("<!-- image -->", " ")
    md = re.sub(r"[#>*`|]+", " ", md)          # markdown artifacts
    md = re.sub(r"[^\S\n]+", " ", md).strip()
    if len(re.findall(r"[A-Za-z][A-Za-z'-]{2,}", md)) < 4:   # too little real text
        return ""
    return md


def _iter_shapes(shapes):
    """Yield every shape, recursing into GROUPs (python-pptx skips these)."""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(sh.shapes)
        else:
            yield sh


# --------------------------------------------------------------------------- #
# Text extraction (keeps page / slide numbers)
# --------------------------------------------------------------------------- #
def extract_pdf(path: Path, ocr: bool = False) -> list[tuple[int, str]]:
    """Return [(page_number_1based, text), ...].

    When ocr=True, image-dominant pages (little native text but embedded images)
    are rendered and OCR'd, so slide-style / scanned pages aren't lost.
    """
    out = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if ocr and len(text) < 40 and page.get_images():
                t = _ocr_image(page.get_pixmap(dpi=200).tobytes("png"), "png")
                if t:
                    text = (text + "\n" + t).strip()
            out.append((i, text))
    return out


def extract_pptx(path: Path, ocr: bool = False) -> list[tuple[int, str]]:
    """Return [(slide_number_1based, text), ...].

    Recurses into GROUP shapes (python-pptx ignores them by default, which is
    how slide-7's title + code were being lost), reads text + tables, and when
    ocr=True OCRs picture shapes to recover image-locked text (e.g. code shots).
    """
    out = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for sh in _iter_shapes(slide.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text.strip())
            elif getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            elif ocr and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = sh.image
                if len(img.blob) >= MIN_OCR_IMG_BYTES:
                    t = _ocr_image(img.blob, img.ext)
                    if t:
                        parts.append(t)
        out.append((i, "\n".join(parts).strip()))
    return out


# --------------------------------------------------------------------------- #
# Chunking
# --------------------------------------------------------------------------- #
def chunk_unit(week, file_obj, course_id: int, path: Path, ocr: bool = False) -> list[dict]:
    """Extract + chunk one file into metadata-tagged chunk records.

    Slides -> one chunk per slide (split only if a slide overflows CHUNK_TOKENS).
    PDF    -> per page, recursively split with overlap. Chunking *per page*
              keeps an exact page number on every chunk for citations.
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        pages, loc_type = extract_pdf(path, ocr=ocr), "page"
    elif ext == ".pptx":
        pages, loc_type = extract_pptx(path, ocr=ocr), "slide"
    else:
        return []

    records = []
    for loc, text in pages:
        if not text:
            continue
        for ci, piece in enumerate(_splitter.split_text(text)):
            records.append(
                {
                    "id": f"{file_obj.id}-{loc_type}{loc}-{ci}",
                    "text": piece,
                    "metadata": {
                        "course_id": course_id,
                        "week": week,
                        "source_file": file_obj.display_name,
                        "file_id": file_obj.id,
                        "loc_type": loc_type,
                        "loc": loc,
                        "chunk_index": ci,
                        "n_tokens": len(_enc.encode(piece)),
                    },
                }
            )
    return records


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #
def main() -> None:
    ap = argparse.ArgumentParser(description="Ingest Canvas content into chunks.")
    ap.add_argument("--course-id", type=int, default=2)
    ap.add_argument("--out", type=Path, default=REPO_ROOT / "data" / "chunks.jsonl")
    ap.add_argument("--raw-dir", type=Path, default=REPO_ROOT / "data" / "raw")
    ap.add_argument("--all-files", action="store_true",
                    help="Ignore modules; ingest every PDF/PPTX on the course.")
    ap.add_argument("--ocr", action="store_true",
                    help="OCR image-locked content (docling; slow, offline). Off by default.")
    args = ap.parse_args()

    course = get_canvas().get_course(args.course_id)
    print(f"Course: {course.name}  (ocr={'on' if args.ocr else 'off'})\n")

    all_records: list[dict] = []
    per_file: list[tuple] = []
    for week, f in iter_content_units(course, all_files=args.all_files):
        path = download(f, args.raw_dir)
        recs = chunk_unit(week, f, args.course_id, path, ocr=args.ocr)
        all_records.extend(recs)
        per_file.append((week, f.display_name, len(recs)))
        print(f"  [{week or '-'}] {f.display_name}: {len(recs)} chunks")

    args.out.parent.mkdir(parents=True, exist_ok=True)
    with args.out.open("w") as fh:
        for r in all_records:
            fh.write(json.dumps(r, ensure_ascii=False) + "\n")

    print(f"\n{len(all_records)} chunks from {len(per_file)} files -> {args.out}")
    if all_records:
        sample = all_records[len(all_records) // 2]
        print("\nSample chunk:")
        print("  id:", sample["id"])
        print("  metadata:", json.dumps(sample["metadata"]))
        print("  text:", sample["text"][:300].replace("\n", " ") + " ...")


if __name__ == "__main__":
    main()
