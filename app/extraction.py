"""Extract text from weekly course materials (PDF, PPTX).

Upgraded to recover content the naive extractors were silently dropping:
  * PPTX GROUP shapes are now recursed into (python-pptx does not descend into
    groups by default — grouped titles / text boxes were being lost), and table
    text is included. These use python-pptx only, so they are ALWAYS ON.
  * OCR (ON by default) recovers text locked inside images — code screenshots,
    diagrams — via docling. Lazy-loaded, so the models load on first use. Pass
    ``ocr=False`` to skip it.

PyMuPDF + docling are regular dependencies. OCR only fires on image-dominant PDF
pages and picture shapes >= MIN_OCR_IMG_BYTES, so text-only slides are untouched.
"""

from __future__ import annotations

import re
import tempfile
from pathlib import Path
from typing import Any

from .course_export import attachment_path

MAX_TEXT_CHARS = 100_000
MIN_TEXT_CHARS = 200
MIN_OCR_IMG_BYTES = 15_000  # skip icons/bullets/logos when OCR-ing pictures

SUPPORTED_MATERIAL_SUFFIXES = (".pdf", ".pptx")


def is_supported_material(filename: str) -> bool:
    """Return True if the file type can be used for quiz generation."""
    lower = filename.lower()
    return lower.endswith(SUPPORTED_MATERIAL_SUFFIXES)


# --------------------------------------------------------------------------- #
# OCR (optional, lazy) — recover text locked inside images
# --------------------------------------------------------------------------- #
_ocr_conv = None


def _ocr_image(blob: bytes, ext: str) -> str:
    """OCR one image (bytes) -> cleaned text, or '' if nothing usable."""
    global _ocr_conv
    if _ocr_conv is None:
        from docling.document_converter import DocumentConverter

        _ocr_conv = DocumentConverter()
    with tempfile.NamedTemporaryFile(suffix=f".{ext}") as tmp:
        tmp.write(blob)
        tmp.flush()
        try:
            md = _ocr_conv.convert(tmp.name).document.export_to_markdown()
        except Exception:
            return ""
    md = md.replace("<!-- image -->", " ")
    md = re.sub(r"[#>*`|]+", " ", md)
    md = re.sub(r"[^\S\n]+", " ", md).strip()
    if len(re.findall(r"[A-Za-z][A-Za-z'-]{2,}", md)) < 4:  # too little real text
        return ""
    return md


def _iter_shapes(shapes):
    """Yield every shape, recursing into GROUPs (python-pptx skips these)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(sh.shapes)
        else:
            yield sh


# --------------------------------------------------------------------------- #
# Per-file extraction
# --------------------------------------------------------------------------- #
def extract_pdf_text(path: Path, ocr: bool = True) -> str:
    if not ocr:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n\n".join(
            t.strip() for page in reader.pages if (t := page.extract_text()) and t.strip()
        )

    # OCR path: PyMuPDF for text + render image-dominant pages through OCR.
    import fitz  # PyMuPDF

    parts: list[str] = []
    with fitz.open(str(path)) as doc:
        for page in doc:
            text = page.get_text("text").strip()
            if len(text) < 40 and page.get_images():
                t = _ocr_image(page.get_pixmap(dpi=200).tobytes("png"), "png")
                if t:
                    text = (text + "\n" + t).strip()
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def extract_pptx_text(path: Path, ocr: bool = True) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in _iter_shapes(slide.shapes):  # recurse into GROUPs
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_parts.append(line)
            elif getattr(shape, "has_table", False):  # tables (were skipped)
                for row in shape.table.rows:
                    slide_parts.append(" | ".join(c.text for c in row.cells))
            elif ocr and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # image text
                img = shape.image
                if len(img.blob) >= MIN_OCR_IMG_BYTES:
                    text = _ocr_image(img.blob, img.ext)
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_file_text(path: Path, ocr: bool = True) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path, ocr=ocr)
    if suffix == ".pptx":
        return extract_pptx_text(path, ocr=ocr)
    raise ValueError(f"Unsupported file type: {path.suffix} ({path.name})")


def extract_week_text(export_root: Path, module: dict[str, Any], ocr: bool = True) -> str:
    """Concatenate extracted text from all attachments in a week module."""
    sections: list[str] = []
    for item in module.get("items", []):
        if item.get("type") != "Attachment":
            continue
        path = attachment_path(export_root, item)
        text = extract_file_text(path, ocr=ocr).strip()
        header = f"## {path.name}"
        sections.append(f"{header}\n\n{text}" if text else f"{header}\n\n(no extractable text)")

    if not sections:
        raise ValueError(
            f"No attachments in module {module.get('name')!r}. "
            "Run course import or pick another week."
        )

    combined = "\n\n---\n\n".join(sections)
    if len(combined) > MAX_TEXT_CHARS:
        combined = combined[:MAX_TEXT_CHARS] + "\n\n[truncated]"
    return combined


def validate_week_text(text: str, week_name: str) -> None:
    if len(text.strip()) < MIN_TEXT_CHARS:
        raise ValueError(
            f"Extracted text for {week_name} is too short ({len(text)} chars). "
            "Check source files or extraction."
        )
